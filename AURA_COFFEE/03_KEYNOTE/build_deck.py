#!/usr/bin/env python3
"""
AURA — REGENERATIVE COFFEE
Master keynote generator, built in the visual language of
20260706_NaturalIntelligence.pdf.

Design system lifted from the reference deck:
  · 16:9, full bleed
  · Display  — Helvetica Neue Bold, ALL CAPS, tight leading, tight tracking
  · Eyebrow  — DM Mono, uppercase, letterspaced, top-left
  · Body     — Helvetica Neue Regular
  · Wordmark — serif 'aura'
  · Palette  — black / white, with pink, green, orange, purple accents
  · Devices  — scattered rotated image cards, circle diagram with black
               label boxes, spectrum axis, process-box flows, n-up image
               grids with mono spec lists, half-photo/half-black splits

Run:  python3 build_deck.py
Out:  AURA_Regenerative_Coffee.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ─────────────────────────────────────────────────────────────────────
# SYSTEM
# ─────────────────────────────────────────────────────────────────────

REPO = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
PUB = os.path.join(REPO, "public")
OUT = os.path.join(os.path.dirname(__file__), "AURA_Regenerative_Coffee.pptx")

W, H = Inches(13.333), Inches(7.5)

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PINK = RGBColor(0xFF, 0x1E, 0x63)
GREEN = RGBColor(0x3F, 0xA4, 0x5B)
ORANGE = RGBColor(0xF2, 0x6B, 0x3A)
PURPLE = RGBColor(0x7B, 0x2D, 0x8E)
GREY = RGBColor(0x9A, 0x9A, 0x9A)

DISPLAY = "Helvetica Neue"
MONO = "DM Mono"
BODY = "Helvetica Neue"
SERIF = "Instrument Serif"

M = Inches(0.55)          # page margin, matches reference eyebrow inset

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def img(p):
    return os.path.join(PUB, p)


# ─────────────────────────────────────────────────────────────────────
# PRIMITIVES
# ─────────────────────────────────────────────────────────────────────

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line=0.92, space_after=0):
    """runs = [(string, size, bold, colour, font, tracking_pts)]"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    first = True
    for item in runs:
        txt, size, bold, colour, font = item[:5]
        track = item[5] if len(item) > 5 else None
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line
        p.space_after = Pt(space_after)
        for i, ln in enumerate(txt.split("\n")):
            r = p.add_run() if i == 0 else p.add_run()
            if i > 0:
                r.text = "\n" + ln
            else:
                r.text = ln
            f = r.font
            f.name = font
            f.size = Pt(size)
            f.bold = bold
            f.color.rgb = colour
            if track is not None:
                # character spacing via raw XML
                f._rPr.set("spc", str(int(track * 100)))
    return tb


def cover_image(s, path, x=0, y=0, w=None, h=None, dim=None):
    """Insert an image cropped to fill the given box (object-fit: cover)."""
    w = w or W
    h = h or H
    iw, ih = Image.open(img(path)).size
    box_ar, im_ar = w / h, iw / ih
    if im_ar > box_ar:            # image wider — crop sides
        draw_h, draw_w = h, int(h * im_ar)
    else:                          # image taller — crop top/bottom
        draw_w, draw_h = w, int(w / im_ar)
    pic = s.shapes.add_picture(img(path), int(x - (draw_w - w) / 2),
                               int(y - (draw_h - h) / 2), draw_w, draw_h)
    # clip to the box using picture crop
    l = (draw_w - w) / 2 / draw_w
    t = (draw_h - h) / 2 / draw_h
    pic.crop_left = pic.crop_right = l
    pic.crop_top = pic.crop_bottom = t
    pic.left, pic.top, pic.width, pic.height = int(x), int(y), int(w), int(h)
    if dim:
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = BLACK
        sh.fill.transparency = 1 - dim
        sh.line.fill.background()
        sh.shadow.inherit = False
        # transparency via XML (python-pptx has no setter)
        from pptx.oxml.ns import qn
        solid = sh.fill._xPr.find(qn('a:solidFill'))
        clr = solid.find(qn('a:srgbClr'))
        alpha = clr.makeelement(qn('a:alpha'), {'val': str(int(dim * 100000))})
        clr.append(alpha)
    return pic


def eyebrow(s, txt, colour=BLACK, y=None):
    text(s, M, y or Inches(0.42), Inches(9), Inches(0.3),
         [(txt.upper(), 11, False, colour, MONO, 1.2)])


def box_label(s, x, y, w, h, txt, fill=BLACK, colour=WHITE, size=11):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = txt.upper()
    r.font.name = MONO
    r.font.size = Pt(size)
    r.font.color.rgb = colour
    r.font.bold = False
    return sh


def line_h(s, x, y, w, colour=BLACK, weight=1.0):
    from pptx.util import Emu as E
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, E(int(Pt(weight))))
    ln.fill.solid()
    ln.fill.fore_color.rgb = colour
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def wordmark(s, colour=WHITE, x=None, y=None, size=40):
    text(s, x or M, y or Inches(3.35), Inches(3), Inches(0.9),
         [("aura", size, False, colour, SERIF)])


# ─────────────────────────────────────────────────────────────────────
# SLIDE PATTERNS
# ─────────────────────────────────────────────────────────────────────

def s_cover(image, right_lines, hand=None):
    s = slide(BLACK)
    cover_image(s, image)
    wordmark(s, WHITE)
    if hand:
        text(s, Inches(9.0), Inches(3.28), Inches(3.8), Inches(0.8),
             [(hand, 19, True, WHITE, DISPLAY)], align=PP_ALIGN.RIGHT, line=1.05)
    text(s, Inches(9.0), Inches(4.35), Inches(3.8), Inches(1.4),
         [(right_lines, 14, True, WHITE, DISPLAY)],
         align=PP_ALIGN.RIGHT, line=1.2)
    return s


def s_black_center(headline, subline=None, scatter=None, accent=None):
    """Big white type on black, optional scattered rotated image cards."""
    s = slide(BLACK)
    if scatter:
        spots = [(0.5, 0.9, 2.0, -9), (4.4, 0.15, 1.5, 7), (10.8, 0.2, 1.7, 12),
                 (11.6, 2.4, 1.9, -6), (0.3, 4.6, 2.1, 8), (5.2, 5.6, 1.5, -11),
                 (8.9, 5.9, 1.6, 5), (11.9, 5.2, 1.8, -8), (3.1, 4.2, 1.2, 14)]
        for i, p in enumerate(scatter[:len(spots)]):
            x, y, sz, rot = spots[i]
            pic = cover_image(s, p, Inches(x), Inches(y), Inches(sz), Inches(sz))
            pic.rotation = rot
    n = headline.count("\n") + 1
    size = {1: 76, 2: 66}.get(n, 54)
    head_y = {1: 2.5, 2: 2.05, 3: 1.45}.get(n, 1.45)
    sub_y = {1: 4.5, 2: 4.55, 3: 4.35}.get(n, 4.35)
    text(s, Inches(1.4), Inches(head_y), Inches(10.53), Inches(2.6),
         [(headline.upper(), size, True, WHITE, DISPLAY, -1.5)],
         align=PP_ALIGN.CENTER, line=0.88)
    if subline:
        text(s, Inches(2.6), Inches(sub_y), Inches(8.13), Inches(2.2),
             [(subline.upper(), 11, False, accent or WHITE, MONO, 0.9)],
             align=PP_ALIGN.CENTER, line=1.6)
    return s


def s_act(word, sub, accent):
    """Act divider — one word, a rule, a line."""
    s = slide(BLACK)
    text(s, M, Inches(2.75), Inches(12.2), Inches(1.5),
         [(word.upper(), 108, True, WHITE, DISPLAY, -3)], line=0.85)
    line_h(s, M, Inches(4.32), Inches(4.2), accent, 3)
    text(s, M, Inches(4.66), Inches(9), Inches(0.6),
         [(sub.upper(), 13, False, WHITE, MONO, 1.1)], line=1.3)
    return s


def s_white_hero(eb, headline, body=None, size=54, accent_body=None):
    s = slide(WHITE)
    eyebrow(s, eb)
    text(s, M, Inches(1.0), Inches(12.2), Inches(3.2),
         [(headline.upper(), size, True, BLACK, DISPLAY, -1.6)], line=0.9)
    if body:
        text(s, M, Inches(5.55), Inches(7.6), Inches(1.4),
             [(body, 15, False, accent_body or BLACK, BODY)], line=1.4)
    return s


def s_photo_hero(image, headline, caption=None, dim=0.35, size=58):
    s = slide(BLACK)
    cover_image(s, image, dim=dim)
    text(s, Inches(1.4), Inches(2.3), Inches(10.5), Inches(2.4),
         [(headline.upper(), size, True, WHITE, DISPLAY, -1.6)],
         align=PP_ALIGN.CENTER, line=0.88)
    if caption:
        text(s, Inches(3.2), Inches(5.5), Inches(6.9), Inches(1.2),
             [(caption.upper(), 10.5, False, WHITE, MONO, 1.0)],
             align=PP_ALIGN.CENTER, line=1.6)
    return s


def s_split(image, eb, headline, body_lines, img_left=True, size=44):
    """Half image, half white — the reference's workhorse."""
    s = slide(WHITE)
    half = Inches(6.4)
    ix = 0 if img_left else W - half
    tx = half + Inches(0.75) if img_left else M
    cover_image(s, image, ix, 0, half, H)
    text(s, tx, Inches(0.42), Inches(5.9), Inches(0.3),
         [(eb.upper(), 11, False, BLACK, MONO, 1.2)])
    text(s, tx, Inches(1.0), Inches(5.9), Inches(2.8),
         [(headline.upper(), size, True, BLACK, DISPLAY, -1.4)], line=0.9)
    runs = [(l, 13, False, BLACK, BODY) for l in body_lines]
    text(s, tx, Inches(4.15), Inches(5.7), Inches(3.0), runs,
         line=1.4, space_after=9)
    return s


def s_grid(eb, headline, items, cols=None, tall=False):
    """n-up image grid with a name + mono spec list beneath. items =
    [(image, name, [spec lines])]"""
    s = slide(WHITE)
    eyebrow(s, eb)
    text(s, M, Inches(0.95), Inches(12.2), Inches(1.4),
         [(headline.upper(), 40, True, BLACK, DISPLAY, -1.4)], line=0.9)
    n = cols or len(items)
    gutter = Inches(0.28)
    total = W - 2 * M
    cw = int((total - gutter * (n - 1)) / n)
    ch = Inches(2.35) if not tall else Inches(3.0)
    top = Inches(3.0) if not tall else Inches(2.7)
    for i, (p, name, specs) in enumerate(items[:n]):
        x = M + i * (cw + gutter)
        cover_image(s, p, x, top, cw, ch)
        text(s, x, top + ch + Inches(0.22), cw, Inches(0.4),
             [(name, 17, False, BLACK, BODY)], line=1.0)
        if specs:
            runs = [(sp.upper(), 9.5, False, BLACK, MONO, 0.8) for sp in specs]
            text(s, x, top + ch + Inches(0.62), cw, Inches(1.5), runs,
                 line=1.5, space_after=3)
    return s


def s_circle(eb, centre_lines, nodes, accent_imgs=None):
    """Circle with four black label boxes at the compass points."""
    s = slide(WHITE)
    eyebrow(s, eb)
    d = Inches(5.15)
    cx, cy = W / 2, Inches(4.15)
    ring = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - d / 2), int(cy - d / 2), d, d)
    ring.fill.background()
    ring.line.color.rgb = BLACK
    ring.line.width = Pt(1.25)
    ring.shadow.inherit = False
    bw, bh = Inches(1.8), Inches(0.32)
    pos = [(cx - bw / 2, cy - d / 2 - bh / 2),            # top
           (cx + d / 2 - bw / 2, cy - bh / 2),            # right
           (cx - bw / 2, cy + d / 2 - bh / 2),            # bottom
           (cx - d / 2 - bw / 2, cy - bh / 2)]            # left
    for (x, y), label in zip(pos, nodes):
        box_label(s, int(x), int(y), bw, bh, label, size=10)
    runs = [(l.upper(), 25, True, BLACK, DISPLAY, -0.8) for l in centre_lines]
    text(s, int(cx - Inches(1.55)), int(cy - Inches(0.62)), Inches(3.1),
         Inches(1.4), runs, align=PP_ALIGN.CENTER, line=0.98, space_after=0)
    if accent_imgs:
        cover_image(s, accent_imgs[0], Inches(10.75), Inches(0.85),
                    Inches(2.0), Inches(1.65))
        if len(accent_imgs) > 1:
            cover_image(s, accent_imgs[1], Inches(0.55), Inches(5.05),
                        Inches(2.0), Inches(1.65))
    return s


def s_spectrum(eb, headline, left, right, left_img, right_img):
    s = slide(WHITE)
    eyebrow(s, eb)
    text(s, Inches(3.9), Inches(2.1), Inches(5.5), Inches(2.0),
         [(headline.upper(), 46, True, BLACK, DISPLAY, -1.5)],
         align=PP_ALIGN.CENTER, line=0.88)
    cover_image(s, left_img, M, Inches(2.35), Inches(3.1), Inches(1.55))
    cover_image(s, right_img, W - M - Inches(3.1), Inches(2.35),
                Inches(3.1), Inches(1.55))
    text(s, M, Inches(4.05), Inches(3.1), Inches(0.3),
         [(left[1].upper(), 9.5, False, BLACK, MONO, 0.9)])
    text(s, W - M - Inches(3.1), Inches(4.05), Inches(3.1), Inches(0.3),
         [(right[1].upper(), 9.5, False, BLACK, MONO, 0.9)],
         align=PP_ALIGN.RIGHT)
    line_h(s, M, Inches(4.95), W - 2 * M, BLACK, 1.1)
    text(s, M, Inches(5.35), Inches(5), Inches(0.6),
         [(left[0], 22, False, GREEN, BODY)])
    text(s, W - M - Inches(5), Inches(5.35), Inches(5), Inches(0.6),
         [(right[0], 22, False, ORANGE, BODY)], align=PP_ALIGN.RIGHT)
    return s


def s_flow(eb, headline, boxes, note=None):
    """Process-box row with arrows, as in BUILDING OUR WORLD MODEL."""
    s = slide(WHITE)
    eyebrow(s, eb)
    text(s, M, Inches(0.95), Inches(12.2), Inches(1.2),
         [(headline.upper(), 40, True, BLACK, DISPLAY, -1.4)], line=0.9)
    n = len(boxes)
    gutter = Inches(0.42)
    total = W - 2 * M
    bw = int((total - gutter * (n - 1)) / n)
    bh = Inches(0.92)
    y = Inches(2.9)
    for i, (label, fill) in enumerate(boxes):
        x = M + i * (bw + gutter)
        box_label(s, x, y, bw, bh, label, fill=fill, size=10)
        if i < n - 1:
            ax = x + bw + Inches(0.1)
            line_h(s, int(ax), int(y + bh / 2), int(gutter - Inches(0.2)),
                   BLACK, 1.0)
    # return arrow underneath
    line_h(s, M, Inches(4.35), W - 2 * M, BLACK, 1.0)
    text(s, M, Inches(4.5), Inches(6), Inches(0.3),
         [("OUTCOME FEEDS THE NEXT BATCH", 9.5, False, BLACK, MONO, 0.9)])
    if note:
        runs = [(l, 15, False, BLACK, BODY) for l in note]
        text(s, M, Inches(5.4), Inches(9.5), Inches(1.6), runs,
             line=1.45, space_after=9)
    return s


def s_biglist(image, eb_left, eb_right, lines):
    """Half photo with stacked black labels, half huge list — the
    EVERY TREE. EVERY LEAF. slide."""
    s = slide(WHITE)
    half = Inches(6.05)
    cover_image(s, image, 0, 0, half, H, dim=0.22)
    text(s, Inches(0.35), Inches(0.42), Inches(4), Inches(0.3),
         [(eb_left.upper(), 10.5, False, WHITE, MONO, 1.2)])
    text(s, half + Inches(0.55), Inches(0.42), Inches(6.2), Inches(0.3),
         [(eb_right.upper(), 10.5, False, BLACK, MONO, 1.1)])
    size = 27 if len(lines) <= 12 else 23
    runs = [(l.upper(), size, True, BLACK, DISPLAY, -0.9) for l in lines]
    text(s, half + Inches(0.85), Inches(1.15), Inches(6.2), Inches(5.9),
         runs, line=1.0, space_after=0)
    return s


def s_record(eb, headline, record_lines, note_lines):
    """The signed-event slide — mono block on a black plate."""
    s = slide(WHITE)
    eyebrow(s, eb)
    text(s, M, Inches(0.95), Inches(12.2), Inches(1.2),
         [(headline.upper(), 40, True, BLACK, DISPLAY, -1.4)], line=0.9)
    plate = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(2.55),
                               Inches(7.0), Inches(2.5))
    plate.fill.solid()
    plate.fill.fore_color.rgb = BLACK
    plate.line.fill.background()
    plate.shadow.inherit = False
    runs = [(l, 15, False, WHITE, MONO, 0.6) for l in record_lines]
    text(s, M + Inches(0.5), Inches(3.0), Inches(6.2), Inches(1.8), runs,
         line=1.6)
    runs2 = [(l, 15, False, BLACK, BODY) for l in note_lines]
    text(s, Inches(8.2), Inches(2.6), Inches(4.6), Inches(3.2), runs2,
         line=1.45, space_after=12)
    return s


def s_end(image):
    s = slide(BLACK)
    cover_image(s, image, dim=0.45)
    wordmark(s, WHITE)
    text(s, Inches(9.0), Inches(3.42), Inches(3.8), Inches(1.2),
         [("REGENERATIVE\nCOFFEE", 14, True, WHITE, DISPLAY)],
         align=PP_ALIGN.RIGHT, line=1.2)
    text(s, Inches(9.0), Inches(4.5), Inches(3.8), Inches(0.4),
         [("ONE REMARKABLE CIRCLE.", 10, False, WHITE, MONO, 0.9)],
         align=PP_ALIGN.RIGHT)
    return s


# ─────────────────────────────────────────────────────────────────────
# THE DECK
# ─────────────────────────────────────────────────────────────────────

# 01 — Cover
s_cover("ecology/images/aura-forest-floor-seedling.jpg",
        "REGENERATIVE\nCOFFEE",
        hand="MISFITS AND\nMONASTICS.")

# 02 — The question
s_black_center(
    "Regenerative\nCoffee",
    "IS THAT A THING?",
    scatter=["journals/coffee/aura-liquid-gold.jpg",
             "herd/images/aura-calf2.jpg",
             "circular/images/aura-mushroom.jpg",
             "journals/fermentation/aura-coffee-flowers.jpg",
             "circular/images/aura-cpp.jpg",
             "shade/images/aura-sunlit-moss-ground.jpg",
             "journals/coffee/aura-banana-wash.jpg",
             "circular/images/aura-microbe.jpg",
             "ecology/images/aura-lichen-on-bark.jpg"],
    accent=PINK)

# 03 — The definition
s_white_hero("THE DEFINITION",
             "Coffee that leaves\nthe land better\nthan it found it.",
             "Not less bad. Better. Nine layers of practice and measurement, "
             "from Vedic-rooted ferments to Aura's own pattern-reading "
             "software, all pointed at that one sentence. Every other claim "
             "in this deck is evidence for it.")

# 04 — The turn
s_photo_hero("circular/images/aura-dung.jpg",
             "Coffee doesn't\nbegin with coffee.",
             "IT BEGINS ABOUT A METRE UNDERGROUND.")

# 05 — The category
s_spectrum("WHAT WE SEE EMERGING",
           "The only claim\nyou can check",
           ("Organic", "SAYS WHAT YOU DIDN'T DO"),
           ("Regenerative", "SAYS WHERE THE GROUND IS GOING"),
           "journals/coffee/aura-coffee-grading.jpg",
           "circular/images/aura-biodynamic.jpg")

# ── ACT ONE ──────────────────────────────────────────────────────────
s_act("Biodynamic", "Grown in a closed loop.", GREEN)

# 07 — Scale
s_white_hero("THE ESTATE",
             "One estate.\nOne harvest\na year.",
             "Aura Estate — Mudigere, Chikmagalur, Karnataka. "
             "150 acres, 100 of them in coffee. Small is the strategy.")

# 08 — The facts
s_photo_hero("journals/land/aura-mudigere-panorama.jpg",
             "150 acres.\n3,600 feet.\n35,000 trees.",
             "RED LATERITE, pH 6.0–6.5  ·  40–100 IN OF RAIN  ·  14–30 °C\n"
             "ARABICA SLN.9 AND SLN.795  ·  WESTERN GHATS, KARNATAKA",
             dim=0.45, size=54)

# 09 — The herd
s_split("herd/images/aura-relationship2.jpg",
        "THE ENGINE",
        "Our best\nfarmers have\nfour legs.",
        ["Fifty-two Malnad Gidda — indigenous to these hills, grazing them "
         "for centuries before anyone thought to write it down.",
         "They eat the same 150 acres the coffee grows in. What they take "
         "from a block, they give back to it.",
         "Almost everything this soil is fed begins in their bodies.",
         "Alongside Jeevamrit, the estate keeps older ferments too — "
         "coconut milk, curd, buttermilk — the same cow-centred method "
         "India was using long before anyone called it biodynamic."])

# 10 — The herd, measured
s_grid("THE ENGINE", "Fifty-two, on their own ground",
       [("herd/images/aura-grassland2.jpg", "52",
         ["Malnad Gidda", "Indigenous Karnataka breed",
          "Estate-bred, line kept pure"]),
        ("herd/images/aura-calf2.jpg", "Daily",
         ["Health, milk, urine per cow", "Dung logged per herd",
          "One passport per animal"]),
        ("circular/images/aura-microbe.jpg", "Lab-tested",
         ["Every preparation read", "Before and after",
          "Nothing untested reaches soil"])])

# 11 — The circle
s_circle("THE MECHANISM",
         ["One remarkable", "circle"],
         ["1. GRASS", "2. THE HERD", "3. PREPARATIONS", "4. THE SOIL"],
         accent_imgs=["journals/fermentation/aura-coffee-flowers.jpg",
                      "journals/biodynamic/aura-young-calves.jpg"])

# 12 — Made by hand
s_split("circular/images/aura-jeevamrut.jpg",
        "THE MECHANISM",
        "Nothing comes in.\nNothing leaves.",
        ["At first light the dung is gathered fresh, weighed, and split three "
         "ways — to the compost pits, to the microbial brews, to the horns.",
         "154,000 litres of Jeevamrit a year. Each barrel stirred by hand "
         "about forty-five minutes a day, one direction then the other.",
         "2,420 kg of Cow Pat Pit a year across fourteen numbered pits, "
         "matured ninety days and dug in ball by ball at a named root.",
         "Coffee husk and prunings, charred on-site into biochar, go back "
         "into the same pits — carbon that would otherwise burn off, "
         "locked into the soil instead.",
         "Aura does not buy fertility. It makes it."],
        img_left=False)

# 13 — What goes in
s_grid("THE PREPARATION", "Three ingredients, one pit",
       [("circular/images/aura-dung.jpg", "Cow dung",
         ["The living culture", "Fresh, from the estate herd",
          "Gathered at dawn"]),
        ("circular/images/aura-cpp.jpg", "The pit",
         ["Fourteen, numbered", "Turned ~45 min a day",
          "Ninety-day cycle"]),
        ("circular/images/aura-saltrock.jpg", "Basalt dust",
         ["Silica, iron, magnesium", "Ground fine",
          "The mineral half"])])

# 14 — The light
s_white_hero("THE METHOD",
             "Measure first.\nPrune later.\nValidate afterwards.",
             "Shade has been managed here by eye for three centuries. "
             "In the 2026 pre-monsoon season, we put a number on it.")

# 15 — The survey
s_grid("BLOCK 3 · BYTON PATTE · 2026", "We measure the light before we cut",
       [("shade/images/aura-dense-forest-interior.jpg", "Zone A",
         ["~33,000 lux", "Dense shade", "Two-thirds filtered"]),
        ("shade/images/aura-tree-canopy-lookup.jpg", "Zone B",
         ["~62,000 lux", "Medium shade", "Inside the CCRI target"]),
        ("shade/images/aura-sunlit-moss-ground.jpg", "Zone C",
         ["~82,000 lux", "Open canopy", "Above target — cut answers zone"])])

# 16 — The canopy
s_photo_hero("shade/images/aura-misty-mountain-forest.jpg",
             "Four stories.\nOne canopy.",
             "SILVER OAK, FIG, JACKFRUIT  ·  ARECA AND HARDWOOD\n"
             "COFFEE IN THE UNDERSTORY  ·  PEPPER, CARDAMOM, COVER CROP\n\n"
             "FIFTY LUX READINGS AN ACRE BEFORE A BRANCH COMES DOWN",
             dim=0.42)

# ── ACT TWO ──────────────────────────────────────────────────────────
s_act("Transparent", "Written down as it happens.", PINK)

# 18 — Everything, measured
s_biglist("journals/land/aura-canopy-noon-2.jpg",
          "ESTATE DATA",
          "NINE STREAMS, EACH ON ITS OWN CLOCK",
          ["Every ferment.", "Every batch.", "Every block.", "Every tree.",
           "Every cow.", "Every spray.", "Every rainfall.", "Every morning.",
           "Every drying day.", "Every soil test.", "Every correction.",
           "Every mistake.", "Every season."])

# 19 — One record
s_record("THE UNIT OF RECORD",
         "One record, in full",
         ["BD 501 · Block 07 · 06:14",
          "waning moon · humidity 78%",
          "by Raju · dung batch G-03"],
         ["The minute. The moon. The hand. The source.",
          "Machines write their own readings. Every human act is signed "
          "from the field, by the person who did the work.",
          "That row is the unit the whole system is built to protect."])

# 20 — The rule
s_white_hero("THE RULE",
             "No untested\nmaterial touches\nthe soil.",
             "Every batch is read in the estate's own lab before and after — "
             "pH, conductivity, microbial colony counts, Trichoderma, "
             "Pseudomonas, Bacillus. A leaf is diagnosed before it is "
             "treated, never the other way round. A batch that comes back "
             "wrong is corrected or held. Never sprayed.")

# 21 — The testing loop
s_flow("THE TESTING LOOP", "Two circles, turning together",
       [("COLLECT AT DAWN", BLACK), ("BREW / COMPOST", GREY),
        ("TEST THE BATCH", BLACK), ("APPLY TO BLOCK", GREY),
        ("RE-READ SOIL AT 90 DAYS", BLACK)],
       note=["Ninety days after a preparation goes on, the block it fed is "
             "retested and cross-referenced against the batch.",
             "Did the microbes take? Did the carbon rise? Did the block "
             "move? Observe, understand, intervene, measure, learn, "
             "regenerate — the same loop, run on every input."])

# 22 — The admission
s_black_center("And here is\nwhat we can't\nprove.",
               "WE CARRY A FRACTION OF A FERMENTING BATCH FORWARD INTO FRESH\n"
               "CHERRY, THE WAY A SHERRY CELLAR CARRIES ITS CULTURE. CAN A\n"
               "COFFEE FERMENT HOLD A HOUSE CULTURE FROM ONE SEASON TO THE\n"
               "NEXT? NOBODY HAS SHOWN IT. WE ARE TESTING IT, YEAR BY YEAR.\n\n"
               "WE'D RATHER TELL YOU THAT THAN SELL YOU THE METAPHOR.",
               accent=PINK)

# 23 — Certification
s_white_hero("SAID BEFORE YOU ASK",
             "We are not\ncertified organic\nor biodynamic.",
             "We practise both, and publish the record instead. The 32-acre "
             "tea block is in organic transition, targeting 2027.")

# 24 — Traceability
s_split("journals/coffee/aura-coffee-grading.jpg",
        "COMPLETE TRANSPARENCY",
        "Animal → preparation\n→ soil → cherry\n→ cup.",
        ["Every lot carries its own file: harvest weight, immature removed, "
         "floats removed, Brix, water TDS, ferment hours, drying days.",
         "pH every fifteen minutes through the ferment. Temperature three "
         "times daily. Screen grading and defect analysis per SCA protocol.",
         "Held in plain JSON and CSV — files that will still open in a decade."])

# ── ACT THREE ────────────────────────────────────────────────────────
s_act("Flavourful", "Six lots. One harvest. No two cups the same.", ORANGE)

# 26 — The discipline
s_white_hero("THE DISCIPLINE",
             "Only fully ripe\ncherries are plucked.",
             "Minimum 95% ripeness. Every float removed. Ferment stopped at "
             "pH 4.2 — by the number, not by nose.")

# 27 — The turn
s_photo_hero("journals/coffee/aura-dry-osmosis.jpg",
             "The same cherry,\npicked the same morning,\nbecomes six coffees.",
             "TERROIR ALONE DOES NOT EXPLAIN THE SPREAD BETWEEN THEM.\n"
             "TERROIR EXPRESSED THROUGH TECHNIQUE DOES.",
             dim=0.52, size=48)

# 28 — The six
s_grid("ONE HARVEST", "Six lots",
       [("journals/coffee/aura-anaerobic-natural.jpg", "001 Anaerobic Natural",
         ["48 hr sealed", "The cherry's own microbes"]),
        ("journals/coffee/aura-dry-osmosis.jpg", "002 Dry Osmosis",
         ["Dried to 45% first", "Blueberry, fig, raisin"]),
        ("journals/coffee/aura-liquid-gold.jpg", "003 Red Honey",
         ["\"Liquid Gold\"", "Honey, caramel, nut"])], cols=3)

s_grid("ONE HARVEST", "Six lots, continued",
       [("journals/coffee/aura-banana-wash.jpg", "004 Banana Wash",
         ["48 hr under banana leaves", "The most Indian lot"]),
        ("journals/coffee/aura-solera-macaceration.jpg", "005 Solera Maceration",
         ["50% carry-forward", "The flagship"]),
        ("journals/coffee/aura-solera-wash.jpg", "006 Solera Wash",
         ["23% Brix — the highest", "The rarest"])], cols=3)

# 30 — The honest footnote
s_split("journals/fermentation/aura-horn.jpg",
        "THE HONEST FOOTNOTE",
        "Altitude and variety\nset the ceiling.",
        ["The ferment only decides how much of that ceiling the cup reaches.",
         "A sealed ferment or a chosen yeast moves a cup by a point or "
         "three — real, measurable, and worth chasing.",
         "Anyone who tells you the process alone makes the coffee is "
         "skipping the harder half of the story."],
        img_left=False)

# 31 — The cup
s_photo_hero("aura-coffee.jpg",
             "The cup is proof\nof the year\nthat made it.",
             "COFFEE IS AGRICULTURAL MEMORY MADE DRINKABLE.",
             dim=0.45)

# 32 — The answer
s_black_center("Yes, it is.",
               "AURA PRESENTS  ·  REGENERATIVE COFFEE\nONE REMARKABLE CIRCLE.",
               accent=WHITE)

# 33 — The invitation
s_grid("THE EXPERIENCE", "Three months. Three completely different estates.",
       [("shade/images/aura-misty-mountain-forest.jpg", "September",
         ["The green months", "Monsoon letting go",
          "The estate at its most biological", "[NEEDS SOURCE]"]),
        ("journals/coffee/aura-coffee-grading.jpg", "December",
         ["First pick", "The wet mill starts running",
          "Six lots begin diverging", "[NEEDS SOURCE]"]),
        ("aura-coffee.jpg", "January",
         ["The table", "Drying beds full",
          "Six coffees, one morning's picking", "[NEEDS SOURCE]"])], cols=3)

# 34 — End card
s_end("ecology/images/aura-forest-floor-seedling.jpg")

prs.save(OUT)
print(f"Saved {OUT}  —  {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
